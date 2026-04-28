from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import os


def create_pptx(slides_data: list) -> str:
    """
    Создает PowerPoint презентацию из данных слайдов
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Настройка стилей
    TITLE_FONT_SIZE = Pt(36)
    CONTENT_FONT_SIZE = Pt(18)

    # Титульный слайд
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI Generated Presentation"
    if hasattr(title, 'text_frame'):
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(99, 102, 241)

    subtitle.text = f"Создано автоматически | {datetime.now().strftime('%d.%m.%Y')}"

    # Слайды с контентом
    for slide_data in slides_data:
        # Используем layout с заголовком и контентом
        content_slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(content_slide_layout)

        # Заголовок слайда
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Без названия")

        # Настройка заголовка
        if hasattr(title_shape, 'text_frame'):
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = TITLE_FONT_SIZE
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(30, 41, 59)

        # Контент слайда
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()

            # Разбиваем контент на пункты
            content = slide_data.get("content", "")

            # Разделяем по разным разделителям
            points = []
            if ";" in content:
                points = [p.strip() for p in content.split(";") if p.strip()]
            elif "\n" in content:
                points = [p.strip() for p in content.split("\n") if p.strip()]
            elif ". " in content:
                points = [p.strip() for p in content.split(". ") if p.strip()]
            else:
                points = [content]

            # Добавляем пункты
            for i, point in enumerate(points):
                if point:
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = point.strip()
                    p.level = 0
                    p.font.size = CONTENT_FONT_SIZE
                    p.font.color.rgb = RGBColor(71, 85, 105)
                    p.space_after = Pt(12)

    # Финальный слайд
    final_slide_layout = prs.slide_layouts[1]
    final_slide = prs.slides.add_slide(final_slide_layout)

    final_title = final_slide.shapes.title
    final_title.text = "Спасибо за внимание!"

    if len(final_slide.placeholders) > 1:
        final_content = final_slide.placeholders[1]
        final_content.text = "Готовы ответить на ваши вопросы\n\nПрезентация создана с помощью AI Presentation Generator"

    # Сохраняем файл
    os.makedirs("generated", exist_ok=True)
    filename = f"generated/presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)

    return filename